"""
Ralph Deep Research - File Generator

Report generation for PDF, Excel, PowerPoint, and CSV formats.
Based on specs/ARCHITECTURE.md and IMPLEMENTATION_PLAN.md Phase 4.4.

Why this module:
- Reporter agent needs to generate professional reports
- Multiple format support for different use cases
- Template-based generation for consistency
- Integration with Jinja2 for dynamic content

Usage:
    from src.tools.file_generator import FileGenerator

    generator = FileGenerator()

    # Generate PDF report
    pdf_path = await generator.generate_pdf(
        aggregated=aggregated_research,
        config=report_config,
        output_path="reports/sess_123/report.pdf"
    )

    # Generate Excel workbook
    excel_path = await generator.generate_excel(
        aggregated=aggregated_research,
        config=report_config,
        output_path="reports/sess_123/data.xlsx"
    )
"""

from __future__ import annotations

import asyncio
import csv
import io
from datetime import datetime
from pathlib import Path
from typing import Any

from jinja2 import Environment, FileSystemLoader, select_autoescape

from src.tools.errors import InvalidInputError, StorageFullError
from src.tools.logging import get_logger
from src.tools.retry import RETRY_CONFIGS, RetryHandler

logger = get_logger(__name__)

# Maximum file sizes (in bytes)
MAX_PDF_SIZE = 20 * 1024 * 1024  # 20 MB
MAX_EXCEL_SIZE = 20 * 1024 * 1024  # 20 MB
MAX_PPTX_SIZE = 20 * 1024 * 1024  # 20 MB


# =============================================================================
# FILE GENERATOR
# =============================================================================


class FileGenerator:
    """
    Generates professional reports in multiple formats.

    Supported formats:
    - PDF: Full report with formatting, charts, tables
    - Excel: Data workbook with multiple sheets
    - PowerPoint: Presentation slides
    - CSV: Raw data export

    Why this design:
    - Template-based for consistent branding
    - Async for non-blocking generation
    - Integrated retry logic for reliability
    - Size validation to prevent storage issues
    """

    def __init__(
        self,
        template_dir: str = "src/templates",
    ) -> None:
        """
        Initialize file generator.

        Args:
            template_dir: Directory containing templates
        """
        self.template_dir = Path(template_dir)
        self.retry_handler = RetryHandler(RETRY_CONFIGS["file_generation"])

        # Initialize Jinja2 environment for HTML templates
        self.jinja_env = Environment(
            loader=FileSystemLoader(str(self.template_dir / "pdf")),
            autoescape=select_autoescape(["html", "xml"]),
        )

        # Register custom filters
        self.jinja_env.filters["format_date"] = self._format_date
        self.jinja_env.filters["format_number"] = self._format_number
        self.jinja_env.filters["format_currency"] = self._format_currency
        self.jinja_env.filters["format_percent"] = self._format_percent

        logger.debug("FileGenerator initialized", template_dir=str(self.template_dir))

    # =========================================================================
    # PDF GENERATION
    # =========================================================================

    async def generate_pdf(
        self,
        aggregated: dict[str, Any],
        config: dict[str, Any],
        output_path: str,
    ) -> str:
        """
        Generate PDF report from aggregated research.

        Args:
            aggregated: AggregatedResearch data
            config: ReportConfig with PDF settings
            output_path: Path for output file

        Returns:
            Path to generated PDF file

        Raises:
            StorageFullError: If generated file exceeds size limit
            InvalidInputError: If template rendering fails
        """
        async def _generate() -> str:
            try:
                # Load PDF template
                template = self.jinja_env.get_template("report.html")
            except Exception:
                # Use default template if custom not found
                template = self._get_default_pdf_template()

            # Get PDF config
            pdf_config = config.get("pdf", {})
            branding = pdf_config.get("branding", {})

            # Prepare template context
            context = {
                "title": aggregated.get("brief", {}).get("goal", "Research Report"),
                "generated_at": datetime.now().isoformat(),
                "executive_summary": aggregated.get("executive_summary", ""),
                "key_insights": aggregated.get("key_insights", []),
                "sections": aggregated.get("sections", []),
                "recommendation": aggregated.get("recommendation", {}),
                "contradictions": aggregated.get("contradictions_found", []),
                "coverage_summary": aggregated.get("coverage_summary", []),
                "metadata": aggregated.get("metadata", {}),
                # Branding
                "logo_url": branding.get("logo_url", ""),
                "primary_color": branding.get("primary_color", "#1a365d"),
                "company_name": branding.get("company_name", ""),
                # Config
                "include_toc": pdf_config.get("include_toc", True),
                "include_charts": pdf_config.get("include_charts", True),
                "language": config.get("language", "en"),
                "style": config.get("style", "formal"),
            }

            # Render HTML
            try:
                html_content = template.render(**context)
            except Exception as e:
                logger.error("Template rendering failed", error=str(e))
                raise InvalidInputError(
                    message=f"Failed to render PDF template: {e}",
                    field="template",
                )

            # Convert HTML to PDF
            pdf_bytes = await self._html_to_pdf(html_content)

            # Validate size
            if len(pdf_bytes) > MAX_PDF_SIZE:
                raise StorageFullError(
                    message="Generated PDF exceeds size limit",
                    storage_type="pdf",
                    current_size_mb=len(pdf_bytes) / (1024 * 1024),
                    limit_mb=MAX_PDF_SIZE / (1024 * 1024),
                )

            # Write file
            output = Path(output_path)
            output.parent.mkdir(parents=True, exist_ok=True)
            output.write_bytes(pdf_bytes)

            logger.info(
                "PDF generated",
                path=str(output),
                size_bytes=len(pdf_bytes),
            )

            return str(output)

        return await self.retry_handler.execute(_generate)

    async def _html_to_pdf(self, html_content: str) -> bytes:
        """
        Convert HTML to PDF using WeasyPrint.

        Falls back to simple HTML-to-text if WeasyPrint unavailable.
        """
        try:
            # Try WeasyPrint first
            from weasyprint import HTML

            # Run in thread pool to avoid blocking
            loop = asyncio.get_event_loop()
            pdf_bytes = await loop.run_in_executor(
                None,
                lambda: HTML(string=html_content).write_pdf()
            )
            return pdf_bytes

        except ImportError:
            logger.warning("WeasyPrint not available, using fallback PDF generation")
            return await self._fallback_pdf_generation(html_content)

    async def _fallback_pdf_generation(self, html_content: str) -> bytes:
        """
        Fallback PDF generation using ReportLab if WeasyPrint unavailable.
        """
        try:
            from reportlab.lib import colors
            from reportlab.lib.pagesizes import letter
            from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
            from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

            buffer = io.BytesIO()
            doc = SimpleDocTemplate(buffer, pagesize=letter)
            styles = getSampleStyleSheet()

            # Create simple PDF from HTML text content
            import re
            text = re.sub(r"<[^>]+>", "", html_content)

            story = []
            for paragraph in text.split("\n\n"):
                if paragraph.strip():
                    story.append(Paragraph(paragraph.strip(), styles["Normal"]))
                    story.append(Spacer(1, 12))

            doc.build(story)
            return buffer.getvalue()

        except ImportError:
            # Ultimate fallback: just return HTML as bytes
            logger.warning("ReportLab not available, returning HTML content")
            return html_content.encode("utf-8")

    def _get_default_pdf_template(self):
        """Create default PDF template if none exists."""
        from jinja2 import Template

        default_html = """
<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <title>{{ title }}</title>
    <style>
        body { font-family: Arial, sans-serif; margin: 40px; color: #333; }
        h1 { color: {{ primary_color }}; border-bottom: 2px solid {{ primary_color }}; padding-bottom: 10px; }
        h2 { color: {{ primary_color }}; margin-top: 30px; }
        h3 { color: #555; }
        .executive-summary { background: #f5f5f5; padding: 20px; border-radius: 5px; margin: 20px 0; }
        .insight { background: #e8f4f8; padding: 15px; margin: 10px 0; border-left: 4px solid {{ primary_color }}; }
        .recommendation { background: #fff3cd; padding: 20px; border-radius: 5px; margin: 20px 0; }
        .section { margin: 30px 0; padding: 20px; border: 1px solid #ddd; border-radius: 5px; }
        .metadata { font-size: 0.9em; color: #666; margin-top: 40px; border-top: 1px solid #ddd; padding-top: 20px; }
        table { width: 100%; border-collapse: collapse; margin: 15px 0; }
        th, td { border: 1px solid #ddd; padding: 10px; text-align: left; }
        th { background: {{ primary_color }}; color: white; }
        .footer { text-align: center; font-size: 0.8em; color: #999; margin-top: 40px; }
    </style>
</head>
<body>
    {% if logo_url %}<img src="{{ logo_url }}" alt="Logo" style="max-height: 60px;">{% endif %}

    <h1>{{ title }}</h1>
    <p>Generated: {{ generated_at }}</p>

    {% if executive_summary %}
    <div class="executive-summary">
        <h2>Executive Summary</h2>
        <p>{{ executive_summary }}</p>
    </div>
    {% endif %}

    {% if key_insights %}
    <h2>Key Insights</h2>
    {% for insight in key_insights %}
    <div class="insight">
        <strong>{{ insight.insight }}</strong>
        {% if insight.supporting_data %}<p>{{ insight.supporting_data }}</p>{% endif %}
    </div>
    {% endfor %}
    {% endif %}

    {% for section in sections %}
    <div class="section">
        <h2>{{ section.title }}</h2>
        <p>{{ section.summary }}</p>
        {% if section.key_points %}
        <h3>Key Points</h3>
        <ul>
        {% for point in section.key_points %}
            <li>{{ point }}</li>
        {% endfor %}
        </ul>
        {% endif %}
    </div>
    {% endfor %}

    {% if recommendation %}
    <div class="recommendation">
        <h2>Recommendation</h2>
        <p><strong>Verdict:</strong> {{ recommendation.verdict }}</p>
        <p><strong>Confidence:</strong> {{ recommendation.confidence }}</p>
        <p>{{ recommendation.reasoning }}</p>
        {% if recommendation.pros %}
        <h3>Pros</h3>
        <ul>{% for pro in recommendation.pros %}<li>{{ pro }}</li>{% endfor %}</ul>
        {% endif %}
        {% if recommendation.cons %}
        <h3>Cons</h3>
        <ul>{% for con in recommendation.cons %}<li>{{ con }}</li>{% endfor %}</ul>
        {% endif %}
    </div>
    {% endif %}

    <div class="metadata">
        <p>Rounds: {{ metadata.total_rounds or 'N/A' }} |
           Data Tasks: {{ metadata.total_data_tasks or 'N/A' }} |
           Research Tasks: {{ metadata.total_research_tasks or 'N/A' }}</p>
    </div>

    <div class="footer">
        {% if company_name %}{{ company_name }} - {% endif %}
        Generated by Ralph Deep Research
    </div>
</body>
</html>
"""
        return Template(default_html)

    # =========================================================================
    # EXCEL GENERATION
    # =========================================================================

    async def generate_excel(
        self,
        aggregated: dict[str, Any],
        config: dict[str, Any],
        output_path: str,
    ) -> str:
        """
        Generate Excel workbook from aggregated research.

        Args:
            aggregated: AggregatedResearch data
            config: ReportConfig with Excel settings
            output_path: Path for output file

        Returns:
            Path to generated Excel file
        """
        async def _generate() -> str:
            try:
                from openpyxl import Workbook
                from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
                from openpyxl.utils import get_column_letter
            except ImportError:
                logger.error("openpyxl not available")
                raise InvalidInputError(
                    message="openpyxl not installed - cannot generate Excel",
                    field="dependency",
                )

            excel_config = config.get("excel", {})
            wb = Workbook()

            # Remove default sheet
            wb.remove(wb.active)

            # Style definitions
            header_fill = PatternFill(start_color="1a365d", end_color="1a365d", fill_type="solid")
            header_font = Font(color="FFFFFF", bold=True)
            thin_border = Border(
                left=Side(style="thin"),
                right=Side(style="thin"),
                top=Side(style="thin"),
                bottom=Side(style="thin"),
            )

            # 1. Summary Sheet
            ws_summary = wb.create_sheet("Summary")
            self._create_summary_sheet(ws_summary, aggregated, header_fill, header_font, thin_border)

            # 2. Key Insights Sheet
            if aggregated.get("key_insights"):
                ws_insights = wb.create_sheet("Key Insights")
                self._create_insights_sheet(ws_insights, aggregated["key_insights"], header_fill, header_font, thin_border)

            # 3. Section Sheets (one per scope item)
            for i, section in enumerate(aggregated.get("sections", []), 1):
                title = section.get("title", f"Section {i}")[:31]  # Excel sheet name limit
                ws_section = wb.create_sheet(title)
                self._create_section_sheet(ws_section, section, header_fill, header_font, thin_border)

            # 4. Data Tables Sheet
            data_tables = self._extract_data_tables(aggregated)
            if data_tables:
                ws_data = wb.create_sheet("Data Tables")
                self._create_data_sheet(ws_data, data_tables, header_fill, header_font, thin_border)

            # 5. Raw Data Sheet (if include_raw_data is True)
            if excel_config.get("include_raw_data", False):
                ws_raw = wb.create_sheet("Raw Data")
                self._create_raw_data_sheet(ws_raw, aggregated)

            # Save workbook
            output = Path(output_path)
            output.parent.mkdir(parents=True, exist_ok=True)

            # Run in thread pool to avoid blocking
            loop = asyncio.get_event_loop()
            await loop.run_in_executor(None, lambda: wb.save(str(output)))

            # Validate size
            file_size = output.stat().st_size
            if file_size > MAX_EXCEL_SIZE:
                output.unlink()  # Delete oversized file
                raise StorageFullError(
                    message="Generated Excel exceeds size limit",
                    storage_type="excel",
                    current_size_mb=file_size / (1024 * 1024),
                    limit_mb=MAX_EXCEL_SIZE / (1024 * 1024),
                )

            logger.info(
                "Excel generated",
                path=str(output),
                size_bytes=file_size,
                sheets=len(wb.sheetnames),
            )

            return str(output)

        return await self.retry_handler.execute(_generate)

    def _create_summary_sheet(self, ws, aggregated, header_fill, header_font, border):
        """Create summary sheet with key metrics."""
        from openpyxl.styles import Alignment

        # Title
        ws["A1"] = "Research Summary"
        ws["A1"].font = header_font
        ws["A1"].fill = header_fill
        ws.merge_cells("A1:D1")

        # Executive Summary
        ws["A3"] = "Executive Summary"
        ws["A3"].font = header_font
        ws["A4"] = aggregated.get("executive_summary", "N/A")
        ws.merge_cells("A4:D4")
        ws["A4"].alignment = Alignment(wrap_text=True)

        # Metadata
        row = 6
        ws[f"A{row}"] = "Metadata"
        ws[f"A{row}"].font = header_font
        row += 1

        metadata = aggregated.get("metadata", {})
        for key, value in [
            ("Total Rounds", metadata.get("total_rounds", "N/A")),
            ("Data Tasks", metadata.get("total_data_tasks", "N/A")),
            ("Research Tasks", metadata.get("total_research_tasks", "N/A")),
            ("Sources", metadata.get("sources_count", "N/A")),
        ]:
            ws[f"A{row}"] = key
            ws[f"B{row}"] = str(value)
            row += 1

        # Auto-size columns
        ws.column_dimensions["A"].width = 20
        ws.column_dimensions["B"].width = 40
        ws.column_dimensions["C"].width = 20
        ws.column_dimensions["D"].width = 40

    def _create_insights_sheet(self, ws, insights, header_fill, header_font, border):
        """Create key insights sheet."""
        # Headers
        headers = ["#", "Insight", "Supporting Data", "Importance"]
        for col, header in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col, value=header)
            cell.fill = header_fill
            cell.font = header_font
            cell.border = border

        # Data
        for row, insight in enumerate(insights, 2):
            ws.cell(row=row, column=1, value=row - 1).border = border
            ws.cell(row=row, column=2, value=insight.get("insight", "")).border = border
            ws.cell(row=row, column=3, value=insight.get("supporting_data", "")).border = border
            ws.cell(row=row, column=4, value=insight.get("importance", "")).border = border

        # Auto-size
        ws.column_dimensions["A"].width = 5
        ws.column_dimensions["B"].width = 50
        ws.column_dimensions["C"].width = 40
        ws.column_dimensions["D"].width = 15

    def _create_section_sheet(self, ws, section, header_fill, header_font, border):
        """Create sheet for a single section."""
        from openpyxl.styles import Alignment

        # Title
        ws["A1"] = section.get("title", "Section")
        ws["A1"].font = header_font
        ws["A1"].fill = header_fill
        ws.merge_cells("A1:C1")

        # Summary
        ws["A3"] = "Summary"
        ws["A3"].font = header_font
        ws["A4"] = section.get("summary", "")
        ws["A4"].alignment = Alignment(wrap_text=True)
        ws.merge_cells("A4:C4")

        # Key Points
        row = 6
        key_points = section.get("key_points", [])
        if key_points:
            ws[f"A{row}"] = "Key Points"
            ws[f"A{row}"].font = header_font
            row += 1
            for point in key_points:
                ws[f"A{row}"] = f"• {point}"
                row += 1

        # Data Highlights
        data_highlights = section.get("data_highlights", {})
        if data_highlights:
            row += 1
            ws[f"A{row}"] = "Data Highlights"
            ws[f"A{row}"].font = header_font
            row += 1
            for key, value in data_highlights.items():
                ws[f"A{row}"] = key
                ws[f"B{row}"] = str(value)
                row += 1

        # Column widths
        ws.column_dimensions["A"].width = 30
        ws.column_dimensions["B"].width = 40
        ws.column_dimensions["C"].width = 30

    def _create_data_sheet(self, ws, tables, header_fill, header_font, border):
        """Create data tables sheet."""
        current_row = 1

        for table in tables:
            # Table name
            ws.cell(row=current_row, column=1, value=table.get("name", "Table")).font = header_font
            current_row += 1

            # Headers
            headers = table.get("headers", [])
            for col, header in enumerate(headers, 1):
                cell = ws.cell(row=current_row, column=col, value=header)
                cell.fill = header_fill
                cell.font = header_font
                cell.border = border
            current_row += 1

            # Rows
            for row_data in table.get("rows", []):
                for col, value in enumerate(row_data, 1):
                    cell = ws.cell(row=current_row, column=col, value=value)
                    cell.border = border
                current_row += 1

            current_row += 2  # Gap between tables

    def _create_raw_data_sheet(self, ws, aggregated):
        """Create raw data sheet with JSON dump."""
        import json
        ws["A1"] = "Raw Aggregated Data (JSON)"
        ws["A2"] = json.dumps(aggregated, indent=2, default=str)

    def _extract_data_tables(self, aggregated: dict[str, Any]) -> list[dict]:
        """Extract data tables from sections."""
        tables = []
        for section in aggregated.get("sections", []):
            section_tables = section.get("data_tables", [])
            tables.extend(section_tables)
        return tables

    # =========================================================================
    # POWERPOINT GENERATION
    # =========================================================================

    async def generate_pptx(
        self,
        aggregated: dict[str, Any],
        config: dict[str, Any],
        output_path: str,
    ) -> str:
        """
        Generate PowerPoint presentation from aggregated research.

        Args:
            aggregated: AggregatedResearch data
            config: ReportConfig with PPTX settings
            output_path: Path for output file

        Returns:
            Path to generated PPTX file
        """
        async def _generate() -> str:
            try:
                from pptx import Presentation
                from pptx.util import Inches, Pt
                from pptx.dml.color import RGBColor
                from pptx.enum.text import PP_ALIGN
            except ImportError:
                logger.error("python-pptx not available")
                raise InvalidInputError(
                    message="python-pptx not installed - cannot generate PowerPoint",
                    field="dependency",
                )

            pptx_config = config.get("pptx", {})
            aspect_ratio = pptx_config.get("aspect_ratio", "16:9")

            # Create presentation with appropriate size
            prs = Presentation()
            if aspect_ratio == "4:3":
                prs.slide_width = Inches(10)
                prs.slide_height = Inches(7.5)
            else:  # 16:9 default
                prs.slide_width = Inches(13.333)
                prs.slide_height = Inches(7.5)

            # Get layouts
            title_layout = prs.slide_layouts[0]  # Title slide
            content_layout = prs.slide_layouts[1]  # Title and content

            # 1. Title Slide
            slide = prs.slides.add_slide(title_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]

            title.text = aggregated.get("brief", {}).get("goal", "Research Report")
            subtitle.text = f"Generated: {datetime.now().strftime('%Y-%m-%d')}"

            # 2. Executive Summary Slide
            slide = prs.slides.add_slide(content_layout)
            slide.shapes.title.text = "Executive Summary"
            body = slide.shapes.placeholders[1]
            tf = body.text_frame
            tf.text = aggregated.get("executive_summary", "No summary available")

            # 3. Key Insights Slides
            insights = aggregated.get("key_insights", [])
            slides_per_section = pptx_config.get("slides_per_section", 2)

            if insights:
                # Split insights across slides if needed
                for i in range(0, len(insights), 4):
                    slide = prs.slides.add_slide(content_layout)
                    slide.shapes.title.text = "Key Insights"
                    body = slide.shapes.placeholders[1]
                    tf = body.text_frame

                    for insight in insights[i:i + 4]:
                        p = tf.add_paragraph() if tf.text else tf.paragraphs[0]
                        p.text = f"• {insight.get('insight', '')}"
                        p.level = 0

            # 4. Section Slides
            for section in aggregated.get("sections", [])[:slides_per_section * 2]:
                slide = prs.slides.add_slide(content_layout)
                slide.shapes.title.text = section.get("title", "Section")

                body = slide.shapes.placeholders[1]
                tf = body.text_frame
                tf.text = section.get("summary", "")[:500]  # Truncate for slide

                # Add key points
                for point in section.get("key_points", [])[:5]:
                    p = tf.add_paragraph()
                    p.text = f"• {point}"
                    p.level = 0

            # 5. Recommendation Slide
            recommendation = aggregated.get("recommendation", {})
            if recommendation:
                slide = prs.slides.add_slide(content_layout)
                slide.shapes.title.text = "Recommendation"

                body = slide.shapes.placeholders[1]
                tf = body.text_frame
                tf.text = f"Verdict: {recommendation.get('verdict', 'N/A')}"

                p = tf.add_paragraph()
                p.text = f"Confidence: {recommendation.get('confidence', 'N/A')}"

                p = tf.add_paragraph()
                p.text = recommendation.get("reasoning", "")[:300]

            # 6. Thank You Slide
            slide = prs.slides.add_slide(title_layout)
            slide.shapes.title.text = "Thank You"
            slide.placeholders[1].text = "Questions?"

            # Save presentation
            output = Path(output_path)
            output.parent.mkdir(parents=True, exist_ok=True)

            loop = asyncio.get_event_loop()
            await loop.run_in_executor(None, lambda: prs.save(str(output)))

            # Validate size
            file_size = output.stat().st_size
            if file_size > MAX_PPTX_SIZE:
                output.unlink()
                raise StorageFullError(
                    message="Generated PPTX exceeds size limit",
                    storage_type="pptx",
                    current_size_mb=file_size / (1024 * 1024),
                    limit_mb=MAX_PPTX_SIZE / (1024 * 1024),
                )

            logger.info(
                "PPTX generated",
                path=str(output),
                size_bytes=file_size,
                slides=len(prs.slides),
            )

            return str(output)

        return await self.retry_handler.execute(_generate)

    # =========================================================================
    # CSV GENERATION
    # =========================================================================

    async def generate_csv(
        self,
        aggregated: dict[str, Any],
        config: dict[str, Any],
        output_path: str,
    ) -> str:
        """
        Generate CSV file(s) from aggregated research data.

        Args:
            aggregated: AggregatedResearch data
            config: ReportConfig with CSV settings
            output_path: Path for output file

        Returns:
            Path to generated CSV file
        """
        async def _generate() -> str:
            csv_config = config.get("csv", {})
            delimiter = csv_config.get("delimiter", ",")
            encoding = csv_config.get("encoding", "utf-8")
            include_headers = csv_config.get("include_headers", True)

            output = Path(output_path)
            output.parent.mkdir(parents=True, exist_ok=True)

            # Collect all data tables
            tables = self._extract_data_tables(aggregated)

            if not tables:
                # Create basic CSV from key insights
                tables = [{
                    "name": "Key Insights",
                    "headers": ["#", "Insight", "Supporting Data"],
                    "rows": [
                        [i + 1, insight.get("insight", ""), insight.get("supporting_data", "")]
                        for i, insight in enumerate(aggregated.get("key_insights", []))
                    ]
                }]

            # Write CSV
            with open(output, "w", newline="", encoding=encoding) as f:
                writer = csv.writer(f, delimiter=delimiter)

                for table in tables:
                    # Table name as comment
                    writer.writerow([f"# {table.get('name', 'Data')}"])

                    # Headers
                    if include_headers:
                        writer.writerow(table.get("headers", []))

                    # Data rows
                    for row in table.get("rows", []):
                        writer.writerow(row)

                    # Blank row between tables
                    writer.writerow([])

            file_size = output.stat().st_size

            logger.info(
                "CSV generated",
                path=str(output),
                size_bytes=file_size,
                tables=len(tables),
            )

            return str(output)

        return await self.retry_handler.execute(_generate)

    # =========================================================================
    # UTILITY METHODS
    # =========================================================================

    @staticmethod
    def _format_date(value: Any, format: str = "%Y-%m-%d") -> str:
        """Jinja2 filter for date formatting."""
        if isinstance(value, datetime):
            return value.strftime(format)
        if isinstance(value, str):
            try:
                dt = datetime.fromisoformat(value.replace("Z", "+00:00"))
                return dt.strftime(format)
            except ValueError:
                return value
        return str(value)

    @staticmethod
    def _format_number(value: Any, decimals: int = 2) -> str:
        """Jinja2 filter for number formatting."""
        try:
            return f"{float(value):,.{decimals}f}"
        except (ValueError, TypeError):
            return str(value)

    @staticmethod
    def _format_currency(value: Any, symbol: str = "$") -> str:
        """Jinja2 filter for currency formatting."""
        try:
            return f"{symbol}{float(value):,.2f}"
        except (ValueError, TypeError):
            return str(value)

    @staticmethod
    def _format_percent(value: Any) -> str:
        """Jinja2 filter for percentage formatting."""
        try:
            return f"{float(value):.1f}%"
        except (ValueError, TypeError):
            return str(value)


# =============================================================================
# FACTORY FUNCTION
# =============================================================================


def get_file_generator(template_dir: str = "src/templates") -> FileGenerator:
    """
    Get file generator instance.

    Args:
        template_dir: Directory containing templates

    Returns:
        Configured FileGenerator instance
    """
    return FileGenerator(template_dir=template_dir)
